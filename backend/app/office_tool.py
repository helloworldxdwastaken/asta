"""Office document generation tool: create .pptx, .docx, and .xlsx files from AI-provided content."""
from __future__ import annotations
import json
import logging
import re
import uuid
from pathlib import Path

logger = logging.getLogger(__name__)

WORKSPACE_DIR = Path(__file__).resolve().parent.parent.parent / "workspace"
OFFICE_OUTPUT_DIR = WORKSPACE_DIR / "office_docs"


def _sanitize_filename(name: str, ext: str) -> str:
    name = re.sub(r"[^\w\s\-.]", "", name).strip()
    if not name or name == ext:
        name = f"document_{uuid.uuid4().hex[:8]}"
    if not name.endswith(ext):
        name += ext
    return name


# ── PowerPoint ────────────────────────────────────────────────────────────────

def is_pptx_available() -> bool:
    try:
        import pptx  # noqa: F401
        return True
    except ImportError:
        return False


def generate_pptx(slides: list[dict], filename: str = "presentation.pptx", theme: str = "dark") -> str:
    """Generate a .pptx file from a list of slide dicts.

    Each slide dict:
      - title: str
      - content: str | list[str]   (bullet points or paragraph text)
      - notes: str (optional speaker notes)
    Returns the output file path.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return "[Error: python-pptx is not installed. Run: pip install python-pptx]"

    OFFICE_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    filename = _sanitize_filename(filename, ".pptx")

    # Theme colours
    is_dark = theme.lower() != "light"
    BG     = RGBColor(0x0F, 0x0F, 0x1A) if is_dark else RGBColor(0xFF, 0xFF, 0xFF)
    TITLE  = RGBColor(0xFF, 0xFF, 0xFF) if is_dark else RGBColor(0x11, 0x11, 0x11)
    BODY   = RGBColor(0xCC, 0xCC, 0xDD) if is_dark else RGBColor(0x33, 0x33, 0x33)
    ACCENT = RGBColor(0xFF, 0x6B, 0x2C)  # Asta orange

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG

        # Accent bar (left edge)
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0),
            Inches(0.12), Inches(7.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        title_text = slide_data.get("title", "")
        content = slide_data.get("content", "")
        notes_text = slide_data.get("notes", "")

        # Title box
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12.5), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = TITLE

        # Divider line
        from pptx.util import Emu as _Emu
        line = slide.shapes.add_shape(
            1,
            Inches(0.4), Inches(1.55),
            Inches(12.5), _Emu(18000),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

        # Content box
        content_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.75), Inches(12.5), Inches(5.4))
        tf2 = content_box.text_frame
        tf2.word_wrap = True

        bullets = content if isinstance(content, list) else [ln for ln in str(content).split("\n") if ln.strip()]
        first = True
        for bullet in bullets:
            p2 = tf2.paragraphs[0] if first else tf2.add_paragraph()
            first = False
            # Strip leading "-" or "•"
            clean = re.sub(r"^[\-•*]\s*", "", bullet.strip())
            p2.text = clean
            p2.font.size = Pt(20)
            p2.font.color.rgb = BODY
            p2.space_before = Pt(4)

        # Speaker notes
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    out_path = OFFICE_OUTPUT_DIR / filename
    prs.save(str(out_path))
    logger.info("Generated PPTX: %s", out_path)
    return str(out_path)


def parse_pptx_tool_args(raw: str) -> dict:
    raw = raw.strip()
    if raw.startswith("{"):
        try:
            return json.loads(raw)
        except json.JSONDecodeError:
            pass
    m = re.search(r"\{[\s\S]*\}", raw)
    if m:
        try:
            return json.loads(m.group())
        except json.JSONDecodeError:
            pass
    return {}


def get_pptx_tool_openai_def() -> list[dict]:
    return [
        {
            "type": "function",
            "function": {
                "name": "generate_pptx",
                "description": (
                    "Create a PowerPoint (.pptx) presentation file. "
                    "Use when the user asks for a presentation, slides, or a deck. "
                    "Build rich slides: each slide should have a clear title and 3-6 bullet points or a paragraph. "
                    "Include speaker notes for complex slides. "
                    "After generating, provide the download link from the result."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {
                            "type": "string",
                            "description": "Output filename (e.g. 'marketing-deck.pptx'). No path.",
                        },
                        "theme": {
                            "type": "string",
                            "enum": ["dark", "light"],
                            "description": "Slide theme. Default: dark.",
                        },
                        "slides": {
                            "type": "array",
                            "description": "List of slide objects.",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "title":   {"type": "string"},
                                    "content": {
                                        "oneOf": [
                                            {"type": "string"},
                                            {"type": "array", "items": {"type": "string"}},
                                        ],
                                        "description": "Bullet list (array) or paragraph (string).",
                                    },
                                    "notes": {"type": "string", "description": "Speaker notes (optional)."},
                                },
                                "required": ["title", "content"],
                            },
                        },
                    },
                    "required": ["slides"],
                },
            },
        }
    ]


# ── Word Document ─────────────────────────────────────────────────────────────

def is_docx_available() -> bool:
    try:
        import docx  # noqa: F401
        return True
    except ImportError:
        return False


def generate_docx(content: str, filename: str = "document.docx", title: str | None = None) -> str:
    """Generate a .docx file from markdown-like content.
    Returns the output file path.
    """
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return "[Error: python-docx is not installed. Run: pip install python-docx]"

    OFFICE_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    filename = _sanitize_filename(filename, ".docx")

    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin    = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin   = Inches(1.25)
        section.right_margin  = Inches(1.25)

    # Title
    if title:
        h = doc.add_heading(title, level=0)
        h.runs[0].font.color.rgb = RGBColor(0xFF, 0x6B, 0x2C)

    # Parse content line by line
    for line in content.split("\n"):
        stripped = line.rstrip()

        # Headings
        hm = re.match(r"^(#{1,4})\s+(.*)", stripped)
        if hm:
            level = len(hm.group(1))
            doc.add_heading(hm.group(2), level=level)
            continue

        # Bullet points
        bm = re.match(r"^[-*+]\s+(.*)", stripped)
        if bm:
            p = doc.add_paragraph(bm.group(1), style="List Bullet")
            p.paragraph_format.space_after = Pt(2)
            continue

        # Numbered list
        nm = re.match(r"^\d+\.\s+(.*)", stripped)
        if nm:
            p = doc.add_paragraph(nm.group(1), style="List Number")
            p.paragraph_format.space_after = Pt(2)
            continue

        # Horizontal rule → page-like separator
        if re.match(r"^(-{3,}|\*{3,})$", stripped):
            doc.add_paragraph("─" * 60)
            continue

        # Empty line → paragraph break
        if not stripped:
            doc.add_paragraph("")
            continue

        # Regular paragraph — handle inline bold/italic
        p = doc.add_paragraph()
        _add_inline_runs(p, stripped)
        p.paragraph_format.space_after = Pt(4)

    out_path = OFFICE_OUTPUT_DIR / filename
    doc.save(str(out_path))
    logger.info("Generated DOCX: %s", out_path)
    return str(out_path)


def _add_inline_runs(paragraph, text: str) -> None:
    """Parse **bold**, *italic*, and plain text into paragraph runs."""
    pattern = re.compile(r"(\*\*(.+?)\*\*|\*(.+?)\*|__(.+?)__|_(.+?)_)")
    last = 0
    for m in pattern.finditer(text):
        if m.start() > last:
            paragraph.add_run(text[last:m.start()])
        raw = m.group(0)
        if raw.startswith("**") or raw.startswith("__"):
            run = paragraph.add_run(m.group(2) or m.group(4))
            run.bold = True
        else:
            run = paragraph.add_run(m.group(3) or m.group(5))
            run.italic = True
        last = m.end()
    if last < len(text):
        paragraph.add_run(text[last:])


def parse_docx_tool_args(raw: str) -> dict:
    raw = raw.strip()
    if raw.startswith("{"):
        try:
            return json.loads(raw)
        except json.JSONDecodeError:
            pass
    m = re.search(r"\{[\s\S]*\}", raw)
    if m:
        try:
            return json.loads(m.group())
        except json.JSONDecodeError:
            pass
    return {"content": raw}


def get_docx_tool_openai_def() -> list[dict]:
    return [
        {
            "type": "function",
            "function": {
                "name": "generate_docx",
                "description": (
                    "Create a Word (.docx) document file. "
                    "Use when the user asks for a report, letter, essay, contract, or any Word document. "
                    "Content supports markdown formatting: # headings, **bold**, *italic*, - bullet lists, 1. numbered lists. "
                    "After generating, provide the download link from the result."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {
                            "type": "string",
                            "description": "Output filename (e.g. 'report.docx'). No path.",
                        },
                        "title": {
                            "type": "string",
                            "description": "Document title shown at the top (optional).",
                        },
                        "content": {
                            "type": "string",
                            "description": "Full document content in markdown format.",
                        },
                    },
                    "required": ["content"],
                },
            },
        }
    ]


# ── Excel Spreadsheet ──────────────────────────────────────────────────────────

def is_xlsx_available() -> bool:
    try:
        import openpyxl  # noqa: F401
        return True
    except ImportError:
        return False


def generate_xlsx(sheets: list[dict], filename: str = "spreadsheet.xlsx", title: str | None = None) -> str:
    """Generate a .xlsx file from a list of sheet dicts.

    Each sheet dict:
      - name: str (sheet tab name)
      - headers: list[str] (column headers)
      - rows: list[list[Any]] (data rows)
      - column_widths: list[int] (optional, per-column widths)
    Returns the output file path.
    """
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
    except ImportError:
        return "[Error: openpyxl is not installed. Run: pip install openpyxl]"

    OFFICE_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    filename = _sanitize_filename(filename, ".xlsx")

    ACCENT_HEX = "FF6B2C"  # Asta orange
    HEADER_BG  = "1A1A2E"  # dark navy
    ALT_ROW_BG = "F5F5FA"  # light lavender tint

    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # remove default sheet

    for sheet_data in sheets:
        sheet_name = str(sheet_data.get("name", "Sheet"))[:31]  # Excel max 31 chars
        ws = wb.create_sheet(title=sheet_name)

        headers: list = sheet_data.get("headers", [])
        rows: list = sheet_data.get("rows", [])
        col_widths: list = sheet_data.get("column_widths", [])

        # Header row
        header_font  = Font(bold=True, color="FFFFFF", size=11)
        header_fill  = PatternFill("solid", fgColor=HEADER_BG)
        header_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
        thin_side    = Side(style="thin", color="DDDDDD")
        thin_border  = Border(left=thin_side, right=thin_side, bottom=thin_side)

        for col_idx, header in enumerate(headers, start=1):
            cell = ws.cell(row=1, column=col_idx, value=str(header))
            cell.font   = header_font
            cell.fill   = header_fill
            cell.alignment = header_align
            cell.border = thin_border

        # Data rows
        alt_fill = PatternFill("solid", fgColor=ALT_ROW_BG)
        data_align = Alignment(vertical="center", wrap_text=True)

        for row_idx, row in enumerate(rows, start=2):
            is_alt = (row_idx % 2 == 0)
            for col_idx, value in enumerate(row, start=1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                cell.alignment = data_align
                cell.border = thin_border
                if is_alt:
                    cell.fill = alt_fill

        # Column widths
        for col_idx, header in enumerate(headers, start=1):
            col_letter = get_column_letter(col_idx)
            if col_idx <= len(col_widths) and col_widths[col_idx - 1]:
                ws.column_dimensions[col_letter].width = col_widths[col_idx - 1]
            else:
                # Auto-size: use max of header length and sample data
                max_len = len(str(header))
                for row in rows[:50]:
                    if col_idx - 1 < len(row):
                        max_len = max(max_len, len(str(row[col_idx - 1])))
                ws.column_dimensions[col_letter].width = min(max_len + 4, 40)

        # Freeze header row
        ws.freeze_panes = "A2"

        # Row height for header
        ws.row_dimensions[1].height = 28

    # Add title sheet at front if requested
    if title and wb.sheetnames:
        ws_title = wb.create_sheet(title="Cover", index=0)
        cell = ws_title.cell(row=2, column=2, value=title)
        cell.font = Font(bold=True, size=20, color=ACCENT_HEX)
        cell.alignment = Alignment(horizontal="left", vertical="center")
        ws_title.column_dimensions["B"].width = 40
        ws_title.row_dimensions[2].height = 40

    out_path = OFFICE_OUTPUT_DIR / filename
    wb.save(str(out_path))
    logger.info("Generated XLSX: %s", out_path)
    return str(out_path)


def parse_xlsx_tool_args(raw: str) -> dict:
    raw = raw.strip()
    if raw.startswith("{"):
        try:
            return json.loads(raw)
        except json.JSONDecodeError:
            pass
    m = re.search(r"\{[\s\S]*\}", raw)
    if m:
        try:
            return json.loads(m.group())
        except json.JSONDecodeError:
            pass
    return {}


def get_xlsx_tool_openai_def() -> list[dict]:
    return [
        {
            "type": "function",
            "function": {
                "name": "generate_xlsx",
                "description": (
                    "Create an Excel (.xlsx) spreadsheet file. "
                    "Use when the user asks for a spreadsheet, tracker, table, budget, schedule, or any data in Excel format. "
                    "Supports multiple sheets, styled headers, and auto-sized columns. "
                    "After generating, provide the download link from the result."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {
                            "type": "string",
                            "description": "Output filename (e.g. 'budget.xlsx'). No path.",
                        },
                        "title": {
                            "type": "string",
                            "description": "Optional cover sheet title.",
                        },
                        "sheets": {
                            "type": "array",
                            "description": "List of sheet objects.",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "name": {"type": "string", "description": "Sheet tab name."},
                                    "headers": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                        "description": "Column header labels.",
                                    },
                                    "rows": {
                                        "type": "array",
                                        "items": {
                                            "type": "array",
                                            "items": {},
                                        },
                                        "description": "Data rows. Each row is an array of values matching the headers.",
                                    },
                                    "column_widths": {
                                        "type": "array",
                                        "items": {"type": "number"},
                                        "description": "Optional explicit column widths (characters). Omit for auto-sizing.",
                                    },
                                },
                                "required": ["name", "headers", "rows"],
                            },
                        },
                    },
                    "required": ["sheets"],
                },
            },
        }
    ]
